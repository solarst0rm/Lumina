"""Core document processing helpers."""

from __future__ import annotations

import base64
import os
import re
import shutil
import tempfile
from pathlib import Path

import fitz
from docx import Document
from openai import OpenAI
from PIL import Image, ImageDraw
from pptx import Presentation

from core.config import (
    API_KEY,
    BASE_URL,
    DEFAULT_EXERCISE_FILENAME,
    DEFAULT_SUMMARY_FILENAME,
    DEFAULT_USER_PROMPT,
    FINAL_PROMPT_TEMPLATE,
    MAX_IMAGE_SIDE,
    MODEL_NAME,
    SECURITY_SENSITIVE_KEYWORDS,
)
from utils.exercise_generator import generate_valid_exercises
from utils.render_utils import render_markdown_to_html


client = OpenAI(base_url=BASE_URL, api_key=API_KEY)

DIRECT_IMAGE_EXTENSIONS = (".jpg", ".jpeg", ".png", ".bmp", ".webp", ".gif", ".tif", ".tiff")
PPT_OOXML_EXTENSIONS = (".pptx", ".pptm", ".ppsx", ".ppsm", ".potx", ".potm")
SUPPORTED_UPLOAD_MESSAGE = (
    "暂不支持该文件格式。请上传 PDF、DOCX、PPTX/PPTM/PPSX/PPSM/POTX/POTM，"
    "或 JPG/JPEG/PNG/BMP/WEBP/GIF/TIFF 图片；旧版 .ppt 请先另存为 .pptx。"
)


def resize_image(path: str) -> None:
    """Resize image in place so the longest side stays within the limit."""
    img = Image.open(path)
    width, height = img.size
    scale = MAX_IMAGE_SIDE / max(width, height)
    if scale < 1:
        resized = img.resize((int(width * scale), int(height * scale)))
        resized.save(path)


def image_to_base64(path: str) -> str:
    """Encode an image file as base64."""
    with open(path, "rb") as file_handle:
        return base64.b64encode(file_handle.read()).decode("utf-8")


def clean_plain_text(text: str) -> str:
    """Remove noisy markup while keeping normal paragraphs intact."""
    cleaned = text or ""
    cleaned = re.sub(r"[\\`|]", "", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def save_markdown(content: str, filename: str | Path) -> None:
    """Save Markdown content to disk."""
    markdown_path = Path(filename)
    markdown_path.parent.mkdir(parents=True, exist_ok=True)
    markdown_path.write_text(content, encoding="utf-8")


def _looks_security_sensitive(text: str) -> bool:
    normalized = (text or "").lower()
    return any(keyword in normalized for keyword in SECURITY_SENSITIVE_KEYWORDS)


def _format_model_error(exc: Exception, stage: str) -> str:
    raw_message = str(exc)
    lowered = raw_message.lower()
    if "output data may contain inappropriate content" in lowered:
        return (
            f"模型拦截了{stage}阶段的结果，因为文档看起来像安全攻防或攻击实验材料。"
            "请改为只生成高层次、教学性、非操作性的总结，减少可直接利用的细节。"
        )
    return raw_message


def _render_text_lines_to_image(lines: list[str], output_path: str) -> None:
    image = Image.new("RGB", (960, 720), "white")
    draw = ImageDraw.Draw(image)
    y_offset = 20
    for line in lines:
        if y_offset > 680:
            break
        draw.text((20, y_offset), line[:80], fill="black")
        y_offset += 22
    image.save(output_path)
    resize_image(output_path)


def _normalize_uploaded_image(filepath: str, temp_dir: str) -> str:
    output_path = os.path.join(temp_dir, f"{Path(filepath).stem}_upload.jpg")
    with Image.open(filepath) as image:
        if getattr(image, "is_animated", False):
            try:
                image.seek(0)
            except EOFError:
                pass
        normalized = image.convert("RGB")
        normalized.save(output_path, format="JPEG", quality=92)
    resize_image(output_path)
    return output_path


def file_to_images(filepath: str) -> tuple[list[str], list[str]]:
    """Convert supported files to a list of images."""
    image_paths: list[str] = []
    temp_images: list[str] = []
    filepath_lower = filepath.lower()
    temp_dir = tempfile.mkdtemp(prefix="doc_pages_")

    try:
        if filepath_lower.endswith(DIRECT_IMAGE_EXTENSIONS):
            normalized_image_path = _normalize_uploaded_image(filepath, temp_dir)
            image_paths.append(normalized_image_path)
            temp_images.append(normalized_image_path)

        elif filepath_lower.endswith(".pdf"):
            document = fitz.open(filepath)
            for index in range(len(document)):
                page = document[index]
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                img_path = os.path.join(temp_dir, f"page_{index}.jpg")
                pixmap.save(img_path)
                resize_image(img_path)
                image_paths.append(img_path)
                temp_images.append(img_path)
            document.close()

        elif filepath_lower.endswith(".ppt"):
            raise ValueError("暂不支持旧版 .ppt 文件，请先另存为 .pptx 后再上传。")

        elif filepath_lower.endswith(PPT_OOXML_EXTENSIONS):
            presentation = Presentation(filepath)
            for index, slide in enumerate(presentation.slides):
                slide_lines = [f"第 {index + 1} 页幻灯片"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_lines.append(shape.text.strip())
                img_path = os.path.join(temp_dir, f"slide_{index}.jpg")
                _render_text_lines_to_image(slide_lines, img_path)
                image_paths.append(img_path)
                temp_images.append(img_path)

        elif filepath_lower.endswith(".docx"):
            document = Document(filepath)
            all_lines: list[str] = []
            for paragraph in document.paragraphs:
                paragraph_text = paragraph.text.strip()
                if paragraph_text:
                    all_lines.append(paragraph_text)
            for table in document.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        all_lines.append(row_text)

            if not all_lines:
                all_lines.append("文档中没有可提取的文字内容。")

            lines_per_page = 30
            for page_index in range(0, len(all_lines), lines_per_page):
                img_path = os.path.join(temp_dir, f"docx_{page_index}.jpg")
                _render_text_lines_to_image(all_lines[page_index : page_index + lines_per_page], img_path)
                image_paths.append(img_path)
                temp_images.append(img_path)

        else:
            raise ValueError("暂不支持该文件格式")

        return image_paths, temp_images
    except Exception:
        shutil.rmtree(temp_dir, ignore_errors=True)
        raise


def extract_reference_text(filepath: str) -> str:
    """Extract plain text for topic detection and prompt context."""
    filepath_lower = filepath.lower()
    chunks: list[str] = []

    try:
        if filepath_lower.endswith(".pdf"):
            document = fitz.open(filepath)
            for index in range(len(document)):
                text = document[index].get_text("text").strip()
                if text:
                    chunks.append(text)
            document.close()

        elif filepath_lower.endswith(PPT_OOXML_EXTENSIONS):
            presentation = Presentation(filepath)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text.strip())

        elif filepath_lower.endswith(".docx"):
            document = Document(filepath)
            for paragraph in document.paragraphs:
                if paragraph.text.strip():
                    chunks.append(paragraph.text.strip())
            for table in document.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        chunks.append(row_text)
    except Exception:
        return ""

    return "\n".join(chunks)[:12000]


def call_model(image_paths: list[str], user_prompt: str, reference_text: str = "", safe_mode: bool = False) -> str:
    """Call the multimodal model to generate a summary."""
    if not API_KEY:
        raise RuntimeError("环境变量 xxx_KEY 未配置")

    final_prompt = FINAL_PROMPT_TEMPLATE.format(user_prompt=user_prompt)
    if safe_mode:
        final_prompt += (
            "\n\n补充要求：\n"
            "- 这份材料可能涉及网络安全实验或攻防教学内容。\n"
            "- 只输出高层次、教学性、非操作性的总结。\n"
            "- 不要提供利用步骤、shellcode、payload、命令、代码或攻击指令。\n"
            "- 优先总结概念、术语、学习目标和防御视角。\n"
        )
    if reference_text:
        final_prompt += f"\n\n以下是文档中的文字片段，可作为理解上下文的参考：\n{reference_text[:4000]}"
    content = [{"type": "text", "text": final_prompt}]
    for path in image_paths:
        content.append(
            {
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{image_to_base64(path)}"},
            }
        )

    response = client.chat.completions.create(
        model=MODEL_NAME,
        messages=[
            {
                "role": "system",
                "content": (
                    "你是一名面向学习场景的中文助教。"
                    "如果材料涉及网络安全、漏洞或利用实验，只能做高层次、教学性、非操作性的讲解。"
                ),
            },
            {"role": "user", "content": content},
        ],
    )
    return response.choices[0].message.content or ""


def process_uploaded_file(filepath: str, user_prompt: str) -> dict:
    """Process an uploaded file and generate summary plus exercises."""
    temp_images: list[str] = []
    temp_dir: str | None = None

    try:
        image_paths, temp_images = file_to_images(filepath)
        if temp_images:
            temp_dir = str(Path(temp_images[0]).resolve().parent)

        final_prompt = user_prompt or DEFAULT_USER_PROMPT
        reference_text = extract_reference_text(filepath)
        safe_mode = _looks_security_sensitive(f"{filepath}\n{reference_text}\n{final_prompt}")
        try:
            model_result = call_model(image_paths, final_prompt, reference_text=reference_text, safe_mode=safe_mode)
        except Exception as exc:
            raise RuntimeError(_format_model_error(exc, "总结生成")) from exc
        markdown_content = clean_plain_text(model_result)

        summary_path = Path(DEFAULT_SUMMARY_FILENAME)
        save_markdown(markdown_content, summary_path)
        render_markdown_to_html(summary_path, is_exercise=False)

        _, exercise_content = generate_valid_exercises(summary_path, safe_mode=safe_mode)

        return {
            "success": True,
            "message": "处理完成",
            "summary": markdown_content,
            "summary_path": str(summary_path.resolve()),
            "exercise_path": str(Path(DEFAULT_EXERCISE_FILENAME).resolve()),
            "exercise": exercise_content,
        }
    except TimeoutError:
        return {"success": False, "error": "处理超时"}
    except Exception as exc:
        return {"success": False, "error": str(exc)}
    finally:
        for path in temp_images:
            if os.path.exists(path):
                try:
                    os.remove(path)
                except OSError:
                    pass
        if temp_dir:
            shutil.rmtree(temp_dir, ignore_errors=True)
